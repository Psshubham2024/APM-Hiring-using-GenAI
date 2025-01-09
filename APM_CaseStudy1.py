import streamlit as st
from docx import Document
import PyPDF2
import pptx
import json
import requests
from io import BytesIO
import time
from queue import Queue
from threading import Thread

# Define the API endpoint and access token
API_URL = "https://api.psnext.info/api/chat"
PSCHATACCESSTOKEN = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJVc2VySW5mbyI6eyJpZCI6MzcxMzcsInJvbGVzIjpbImRlZmF1bHQiXSwicGF0aWQiOiI2ZGUxNDUxMy1jYmI3LTQ0NmYtOTM3ZS0xYzhkNTZiNjBhZDUifSwiaWF0IjoxNzM2NDI5NjI3LCJleHAiOjE3MzkwMjE2Mjd9.HaeQwApGKUreK46i-oprwl76y2znYgVQZ0jozaUCmyc"  # Replace with your actual access token

# Rate limiting variables
requests_per_minute = 5  # Set your API rate limit here
rate_limit_delay = 60 / requests_per_minute  # Time to wait between requests in seconds

# Queue to manage requests
request_queue = Queue()
results = {}

# Function to extract text from a Word document
def extract_text_from_word(docx_file):
    doc = Document(docx_file)
    return '\n'.join([para.text for para in doc.paragraphs])

# Function to extract text from a PDF file
def extract_text_from_pdf(pdf_file):
    reader = PyPDF2.PdfReader(pdf_file)
    pdf_text = ""
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        pdf_text += page.extract_text() or ""  # Ensure no None values
    return pdf_text

# Function to extract text from a PPT file
def extract_text_from_ppt(ppt_file):
    ppt = pptx.Presentation(ppt_file)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Function to extract text from different file types
def extract_text_from_file(file):
    if file.type == "application/pdf":
        return extract_text_from_pdf(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_word(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_ppt(file)
    else:
        return None

# Function to compare two sets of texts (questions and solutions)
def compare_question_solution(question_text, solution_text):
    payload = {
    "message": (
        "You are tasked with evaluating candidate submissions for the role of Agile Program Manager. Candidates receive a case study document (uploaded under 'Input: Upload Case Study Document (PDF)') and are required to submit their responses in the form of a PowerPoint presentation (uploaded under 'Upload Solution Documents (PPT)').\n\n"
        "Please assess the quality of their responses based on the following criteria, providing detailed feedback for each. Rate each response on a scale of 10, and provide an overall rating for the entire response on a scale of 50.\n\n"
        "Evaluation Criteria:\n\n"
        "1. **Project Plan**: Thoroughly assess whether a comprehensive project plan/schedule is provided. If so, evaluate how well it incorporates Agile/Scrum principles, including clarity of timeline, sequence of activities, sprint schedules, release milestones, and stages for SIT, UAT, and Go-Live. Provide detailed feedback on each aspect, noting any elements that are missing or require improvement.\n\n"
        "2. **Understanding of Requirements as a Project Manager**: Evaluate the clarity and completeness of the project requirements outlined by the candidate, ensuring they directly address the case study. Note any gaps or missing elements in their understanding.\n\n"
        "3. **Engagement Objectives and Project Requirements**: Assess how well the candidate has articulated the problem statement, high-level requirements, key business drivers, and goals. Additionally, review their understanding of both business and system implications for implementing the proposed solution.\n\n"
        "4. **Type of Engagement**: Review whether the candidate has detailed the type of engagement, including their recommended pricing model. Check if they provided a transparent breakdown of 'price to client,' covering all costs considered in the price proposal.\n\n"
        "5. **Business Context Diagram**: Evaluate the relevance and effectiveness of the business context diagram in supporting the project proposal. Confirm whether it aligns with the case study’s goals and clarifies project scope and relationships.\n\n"
        "6. **Key Metrics/Tools/Reporting**: Assess the inclusion and clarity of key metrics, tools, and reporting methods proposed for effective delivery management. Evaluate whether the candidate has suggested specific dashboards or reporting mechanisms for project monitoring.\n\n"
        "7. **Risk Management**: Check whether the candidate has identified potential project management risks related to cost, scope, and scheduling. Evaluate their risk identification and mitigation skills as demonstrated in their response.\n\n"
        f"Question Document:\n{question_text}\n\nSolution Document:\n{solution_text}"
    ),
    "options": {"model": "gpt35turbo"}
}

    headers = {
        "Authorization": f"Bearer {PSCHATACCESSTOKEN}",
        "Content-Type": "application/json"
    }

    response = requests.post(API_URL, headers=headers, json=payload)

    if response.status_code == 200:
        response_data = response.json()
        messages = response_data.get('data', {}).get('messages', [])
        for message in messages:
            if message.get('role') == 'assistant':
                return message.get('content', 'No content returned from the API.')
        return 'No assistant message found in the API response.'
    else:
        return f"Error: {response.status_code}, {response.text}"

# Worker function to process requests from the queue
def process_requests():
    while True:
        question_text, solution_text, request_id = request_queue.get()
        if request_id is None:  # Stop signal
            break

        result = compare_question_solution(question_text, solution_text)
        results[request_id] = result
        time.sleep(rate_limit_delay)  # Rate limiting delay
        request_queue.task_done()

# Start the worker thread
worker_thread = Thread(target=process_requests)
worker_thread.start()

# Function to generate a Word document of the evaluation feedback
def create_word_report(feedback):
    doc = Document()
    doc.add_heading('Evaluation Feedback Report', level=1)

    for line in feedback.split('\n'):
        doc.add_paragraph(line)

    output = BytesIO()
    doc.save(output)
    output.seek(0)  # Move to the beginning of the BytesIO buffer
    return output

# Main app page
def main_app():
    st.title("Case Study Evaluation Designed for Agile Program Manager")
    st.write("Upload the Case Study & Candidate Response-Get AI Powered Evaluation with Summary & Ratings!")

    uploaded_question_doc = st.file_uploader("Upload Case Study Document (PDF)", type=["pdf"], key="question_doc")
    uploaded_solution_docs = st.file_uploader("Upload Solution Documents (PPT)", type=["pptx"], accept_multiple_files=True, key="solution_docs")

    if uploaded_question_doc and uploaded_solution_docs:
        question_text = extract_text_from_file(uploaded_question_doc)

        # Input Validation: Check for valid question document upload
        if question_text is None:
            st.error("The uploaded question document is not valid. Please upload a valid PDF file.")
            return

        # Confirmation message for valid uploads
        st.success("Document uploaded successfully! Ready for comparison.")

        if st.button("Evaluate", key="compare_documents_button"):
            with st.spinner("Processing..."):
                combined_feedback = ""
                for solution_file in uploaded_solution_docs:
                    solution_text = extract_text_from_file(solution_file)

                    # Input Validation: Check for valid solution document upload
                    if solution_text is None:
                        st.error(f"The uploaded solution document '{solution_file.name}' is not valid. Please upload a valid PowerPoint file.")
                        continue
                    
                    request_id = f"request_{time.time()}"  # Unique ID for this request
                    request_queue.put((question_text, solution_text, request_id))

                # Display feedback in real time
                for solution_file in uploaded_solution_docs:
                    while True:
                        if len(results) > 0:
                            for request_id in list(results.keys()):
                                feedback = results.pop(request_id)  # Get the result and remove from the dictionary
                                st.text_area(f"Comparison Feedback for {solution_file.name}", feedback, height=300, key=f"comparison_feedback_{request_id}")
                                combined_feedback += f"Feedback for {solution_file.name}:\n{feedback}\n\n"  # Combine feedback
                            break
                        time.sleep(1)  # Check every second

                # Generate and allow downloading of the Word report
                if combined_feedback:  # Only create report if there's feedback
                    word_report = create_word_report(combined_feedback)
                    st.download_button("Download Evaluation Report", word_report, "evaluation_report.docx")

# Footer update with black and white text and reduced box size
    st.markdown(
        """
        <style>
        .footer {
            background-color: white;
            color: black;
            text-align: center;
            padding: 5px;
            font-size: 12px;
        }
        </style>
        """, 
        unsafe_allow_html=True
    )
    st.markdown(
        "<div class='footer'><p>Interview Insight Pro | Powered by Gen AI</p></div>", 
        unsafe_allow_html=True
    )

# Run the app
if __name__ == "__main__":
    main_app()
    request_queue.put((None, None, None))  # Stop the worker thread
    worker_thread.join()  # Wait for the thread
